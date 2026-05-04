"""src/processing/parsers.py — Parseurs multi-format."""
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)


def parse_document(path: Path) -> ParsedDocument:
    """Dispatche vers le bon parseur selon l'extension."""
    ext = path.suffix.lower()
    parsers = {
        ".pdf":  _parse_pdf,
        ".txt":  _parse_text,
        ".md":   _parse_text,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pptx": _parse_pptx,
    }
    fn = parsers.get(ext)
    if fn is None:
        raise ValueError(f"Format non supporté: {ext}")
    doc = fn(path)
    doc.metadata.setdefault("source", str(path))
    doc.metadata.setdefault("filename", path.name)
    doc.metadata.setdefault("word_count", len(doc.content.split()))
    return doc


def _parse_text(path: Path) -> ParsedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ParsedDocument(content=text)


def _parse_pdf(path: Path) -> ParsedDocument:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        pages = [page.get_text() for page in doc]
        text = "\n\n".join(pages)
        return ParsedDocument(content=text, metadata={"pages": len(pages)})
    except ImportError:
        logger.warning("PyMuPDF not installed. Run: pip install pymupdf")
        return ParsedDocument(content="", metadata={"error": "pymupdf missing"})
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        return ParsedDocument(content="", metadata={"error": str(e)})


def _parse_docx(path: Path) -> ParsedDocument:
    try:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return ParsedDocument(content=text)
    except ImportError:
        logger.warning("python-docx not installed. Run: pip install python-docx")
        return ParsedDocument(content="", metadata={"error": "python-docx missing"})


def _parse_xlsx(path: Path) -> ParsedDocument:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"=== {sheet} ===\n" + "\n".join(rows))
        return ParsedDocument(content="\n\n".join(parts))
    except ImportError:
        logger.warning("openpyxl not installed. Run: pip install openpyxl")
        return ParsedDocument(content="", metadata={"error": "openpyxl missing"})


def _parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return ParsedDocument(content="\n\n".join(slides), metadata={"slides": len(prs.slides)})
    except ImportError:
        logger.warning("python-pptx not installed. Run: pip install python-pptx")
        return ParsedDocument(content="", metadata={"error": "python-pptx missing"})
